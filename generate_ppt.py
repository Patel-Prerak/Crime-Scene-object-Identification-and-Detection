from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AR Forensics Object Specification and Identification"
    subtitle.text = "AI-Powered Evidence Detection & Immersive Crime Scene Reconstruction\n\nPresenters: Paarth Birla, Dev Lad, Prerak Patel\nGuidance: Prof. Ranjit Kolkar\nNational Forensic Sciences University (NFSU), Goa Campus"

    # Helper to add content slides
    def add_content_slide(title_text, content_text):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        # Access the body text placeholder
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = content_text
        
        # Optional: formatting
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.space_after = Pt(14)

    # --- Slide 2: Introduction ---
    content_intro = (
        "Current Challenge:\n"
        "- Traditional crime scene investigation is time-consuming (manual documentation).\n"
        "- High risk of evidence contamination and human error.\n"
        "- Chain of Custody is difficult to maintain with physical logs.\n\n"
        "Our Solution:\n"
        "- Automated, AI-driven system for real-time evidence identification.\n"
        "- Bridges Computer Vision (CV) and Forensic Science.\n"
        "- Foundation for immersive AR/VR reconstruction."
    )
    add_content_slide("Introduction", content_intro)

    # --- Slide 3: Purpose and Scope ---
    content_purpose = (
        "Purpose:\n"
        "- Automation: Reduce cognitive load (auto-flag weapons, blood).\n"
        "- Integrity: Immutable digital record immediately upon arrival.\n"
        "- Safety: Remote assessment of dangerous scenes.\n\n"
        "Scope:\n"
        "- Current: Web-based analysis tool for static images (2.5D).\n"
        "- Immediate: AR glasses for heads-up display.\n"
        "- Long-term: Virtual Training, Remote Forensics, Courtroom Visualization."
    )
    add_content_slide("Purpose and Scope", content_purpose)

    # --- Slide 4: Literature Survey ---
    content_lit = (
        "Manual Methods:\n"
        "- Grid search, numbered markers. Effective but slow.\n\n"
        "Computer Vision in Forensics:\n"
        "- Redmon et al. (YOLO): Revolutionized real-time detection.\n"
        "- Previous Studies: Mostly single-class (e.g., only guns).\n"
        "- Gap: Lack of integrated systems (Multi-class + 3D).\n\n"
        "Our Contribution:\n"
        "- Dual-model ensemble + Depth estimation for 3D context."
    )
    add_content_slide("Literature Survey", content_lit)

    # --- Slide 5: Methodology ---
    content_method = (
        "1. Ensemble Evidence Detector:\n"
        "- Model A (Standard YOLOv8l): General objects (Person, electronics).\n"
        "- Model B (Custom Forensic): Guns, Knives, Blood Stains.\n\n"
        "2. Detection Logic:\n"
        "- Parallel Inference -> Confidence Filtering (>30%) -> Result Merging.\n\n"
        "3. 3D Reconstruction:\n"
        "- Depth Estimation (Monocular).\n"
        "- Displacement Map extrusion in WebVR (A-Frame)."
    )
    add_content_slide("Methodology", content_method)

    # --- Slide 6: Technology Stack ---
    content_tech = (
        "Core Logic: Python 3.9+\n\n"
        "Computer Vision / AI:\n"
        "- Ultralytics YOLOv8 detection\n"
        "- OpenCV Processing\n"
        "- NumPy & Pandas\n\n"
        "Interfaces:\n"
        "- Streamlit (Web UI - Forensic Theme)\n"
        "- A-Frame (HTML5 VR Rendering)"
    )
    add_content_slide("Technology Stack", content_tech)

    # --- Slide 7: Model Overview ---
    content_model = (
        "YOLOv8 (You Only Look Once - v8):\n"
        "- Selected for speed (real-time) and accuracy.\n"
        "- Anchor-free detection head (better for small objects).\n\n"
        "Training Strategy:\n"
        "- Transfer Learning on Custom Model.\n"
        "- Heavy augmentation (mosaic, scaling) to simulate crime scene conditions."
    )
    add_content_slide("Model Overview", content_model)

    # --- Slide 8: Model Performance & Features ---
    content_perf = (
        "Performance:\n"
        "- Accuracy: >85% mAP.\n"
        "- Latency: < 2 seconds/image.\n"
        "- Reduced False Negatives via Ensemble.\n\n"
        "Key Features:\n"
        "- 'Forensic Light' Theme (High contrast UI).\n"
        "- Real-time 'CRITICAL' alerts for weapons.\n"
        "- Evidence Badges (Color-coded).\n"
        "- VR Mode (Instant 3D conversion)."
    )
    add_content_slide("Model Performance & Features", content_perf)

    # --- Slide 9: Actual Project Analysis (Post-Update) ---
    content_actual = (
        "Ensemble Integration:\n"
        "- Successfully runs two heavy models in parallel without crashing.\n\n"
        "VR Visualization:\n"
        "- Projects detections into 3D space using depth maps.\n"
        "- Users can navigate the scene in-browser.\n"
        "- Labels float at correct depth, not just on screen 2D plane.\n\n"
        "UI Refinement:\n"
        "- Polished 'Glassmorphism' design with purple/white palette."
    )
    add_content_slide("Actual Project Analysis", content_actual)

    # --- Slide 10: Future Scope ---
    content_future = (
        "1. VR Training:\n"
        "- Students explore procedural crime scenes safely.\n\n"
        "2. Court Admissibility:\n"
        "- 'Virtual Jury' - Immersive 3D walkthroughs for context.\n"
        "- Verifying witness lines-of-sight.\n\n"
        "3. Professional Tools:\n"
        "- Remote Expert Consultation (Live 3D stream).\n"
        "- Auto-Report Generation (Chain of Custody Compliant)."
    )
    add_content_slide("Future Scope & Integration", content_future)

    # --- Slide 11: Conclusion ---
    content_conc = (
        "- Demonstrates consumer-grade AI for forensic workflows.\n"
        "- Combines Ensemble Object Detection with VR Visualization.\n"
        "- Adds a new dimension to investigation: Depth and Immersion."
    )
    add_content_slide("Conclusion", content_conc)

    # Save
    output_path = "AR_Forensics_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
